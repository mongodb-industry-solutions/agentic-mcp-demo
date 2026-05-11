#
# Copyright (c) 2026 MongoDB Inc.
# Author: Benjamin Lorenz <benjamin.lorenz@mongodb.com>
#

"""
ACC Export Service — Export MongoDB sales proof points to PowerPoint slides and HTML one-pagers.

Use for: "export to pptx", "generate the slide", "make a PowerPoint", "export to html",
"create the one-pager", "publish the proof point", "deploy it", "save as slide",
"generate the deck". This is a sales content export tool, NOT a network or monitoring tool.

Route ALL of the following here:
- Export:  "export to pptx", "generate the slide", "make a PowerPoint",
           "export to html", "create the one-pager", "save as slide",
           "publish the proof point", "deploy it", "create the file",
           "generate the deck", "produce the one-pager"
- List:    "list exports", "what have we exported", "show generated files"

Does NOT capture proof point content — use acc_proof_point_service for that.
Blocks export if required sections (problem, solution, results, KPIs) are missing.
"""

import datetime
import logging
import os
from pathlib import Path
from pymongo import MongoClient, DESCENDING
from mcp.server.fastmcp import FastMCP

logging.disable(logging.WARNING)

mcp    = FastMCP("acc_export_service")
logger = logging.getLogger("acc_export_service")

mongo_client = MongoClient(os.environ["MONGODB_URI"])
db           = mongo_client["agent_registry"]
proof_points = db["acc_proof_points"]

EXPORTS_DIR = Path(__file__).parent.parent / "exports"
EXPORTS_DIR.mkdir(exist_ok=True)

MDB_GREEN  = "00ED64"
MDB_DARK   = "001E2B"
MDB_GRAY   = "E8EDEB"
MDB_WHITE  = "FFFFFF"


# ─── Helpers ────────────────────────────────────────────────────────────────

def _active() -> dict | None:
    """Most complete non-published proof point, falling back to most recently updated."""
    import datetime as _dt
    drafts = list(proof_points.find({"status": {"$ne": "published"}}))
    if not drafts:
        return None

    def _pct_local(doc):
        prob = doc.get("problem") or {}
        sol  = doc.get("solution") or {}
        res  = doc.get("results") or {}
        checks = [
            bool(doc.get("customer")), bool(doc.get("use_case")),
            bool(prob.get("situation")), bool(prob.get("negative_outcomes")),
            bool(sol.get("what_was_built")), bool(sol.get("mongodb_capabilities")),
            bool(sol.get("why_mongodb")), bool(res.get("kpis")),
        ]
        return sum(checks)

    return max(drafts, key=lambda d: (_pct_local(d), d.get("updated_at") or _dt.datetime.min))


REQUIRED = {
    "customer":                      "customer / company name",
    "use_case":                      "use case",
    "problem.situation":             "problem situation",
    "problem.negative_outcomes":     "negative outcomes",
    "solution.what_was_built":       "solution description",
    "solution.mongodb_capabilities": "MongoDB capabilities",
    "solution.why_mongodb":          "why MongoDB",
    "results.kpis":                  "KPIs",
}


def _missing(doc: dict) -> list[str]:
    prob = doc.get("problem") or {}
    sol  = doc.get("solution") or {}
    res  = doc.get("results") or {}
    checks = {
        "customer":                      bool(doc.get("customer")),
        "use_case":                      bool(doc.get("use_case")),
        "problem.situation":             bool(prob.get("situation")),
        "problem.negative_outcomes":     bool(prob.get("negative_outcomes")),
        "solution.what_was_built":       bool(sol.get("what_was_built")),
        "solution.mongodb_capabilities": bool(sol.get("mongodb_capabilities")),
        "solution.why_mongodb":          bool(sol.get("why_mongodb")),
        "results.kpis":                  bool(res.get("kpis")),
    }
    return [REQUIRED[k] for k, ok in checks.items() if not ok]


def _slug(doc: dict) -> str:
    customer = (doc.get("customer") or "unknown").lower().replace(" ", "-").replace("&", "and")
    pid      = doc["_id"].lower()
    ts       = datetime.datetime.now().strftime("%Y%m%d")
    return f"{pid}-{customer}-{ts}"


# ─── HTML export ────────────────────────────────────────────────────────────

def _export_html_file(doc: dict) -> Path:
    prob = doc.get("problem") or {}
    sol  = doc.get("solution") or {}
    res  = doc.get("results") or {}

    neg_outcomes_html = "".join(
        f'<li>{o}</li>' for o in (prob.get("negative_outcomes") or [])
    )
    caps_html = " &nbsp;·&nbsp; ".join(
        f'<span class="cap">{c}</span>' for c in (sol.get("mongodb_capabilities") or [])
    )
    kpis_html = "".join(
        f'<div class="kpi">{k}</div>' for k in (res.get("kpis") or [])
    )
    outcomes_html = "".join(
        f'<li>{o}</li>' for o in (res.get("outcomes") or [])
    )

    html = f"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="utf-8"/>
<meta name="viewport" content="width=device-width,initial-scale=1"/>
<title>{doc.get('customer','—')} — MongoDB Proof Point</title>
<link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap" rel="stylesheet"/>
<style>
  *, *::before, *::after {{ box-sizing: border-box; margin: 0; padding: 0; }}
  body {{
    font-family: 'Inter', sans-serif;
    background: #f8f9fa; color: #001E2B;
    min-height: 100vh; padding: 48px 32px;
    -webkit-font-smoothing: antialiased;
  }}
  .page {{
    max-width: 960px; margin: 0 auto;
    background: #fff; border-radius: 16px;
    box-shadow: 0 4px 24px rgba(0,30,43,0.10);
    overflow: hidden;
  }}
  .header {{
    background: #001E2B; padding: 36px 48px;
    display: flex; justify-content: space-between; align-items: flex-start;
  }}
  .header-left h1 {{
    font-size: 28px; font-weight: 700; color: #fff;
    letter-spacing: -0.02em; margin-bottom: 6px;
  }}
  .header-left .use-case {{
    font-size: 15px; color: rgba(255,255,255,0.6); font-weight: 400;
  }}
  .header-right {{
    font-family: monospace; font-size: 12px;
    color: rgba(255,255,255,0.35); text-align: right; margin-top: 4px;
  }}
  .mdb-badge {{
    background: #00ED64; color: #001E2B;
    font-size: 11px; font-weight: 700; letter-spacing: 0.08em;
    text-transform: uppercase; padding: 4px 12px; border-radius: 4px;
    display: inline-block; margin-top: 12px;
  }}
  .body {{ padding: 0; display: grid; grid-template-columns: 1fr 1fr 1fr; }}
  .col {{
    padding: 36px 32px; border-right: 1px solid #E8EDEB;
  }}
  .col:last-child {{ border-right: none; }}
  .col-label {{
    font-size: 10px; font-weight: 700; text-transform: uppercase;
    letter-spacing: 0.12em; margin-bottom: 6px;
  }}
  .col.problem  .col-label {{ color: #d64045; }}
  .col.solution .col-label {{ color: #0064FF; }}
  .col.results  .col-label {{ color: #00684A; }}
  .col h2 {{
    font-size: 16px; font-weight: 700; margin-bottom: 16px;
    color: #001E2B; line-height: 1.3;
  }}
  .col p {{ font-size: 13px; color: #3d4f58; line-height: 1.65; margin-bottom: 12px; }}
  .col ul {{ padding-left: 16px; margin-bottom: 12px; }}
  .col ul li {{ font-size: 13px; color: #3d4f58; line-height: 1.7; }}
  .caps {{ margin: 14px 0; display: flex; flex-wrap: wrap; gap: 6px; }}
  .cap {{
    font-size: 11px; font-weight: 600;
    background: #E8F4FE; color: #0064FF;
    border-radius: 4px; padding: 3px 9px;
  }}
  .why-mdb {{
    font-size: 12px; font-style: italic; color: #3d4f58;
    border-left: 3px solid #00ED64; padding-left: 10px; margin-top: 12px;
    line-height: 1.6;
  }}
  .kpi {{
    background: #00ED64; color: #001E2B;
    font-size: 13px; font-weight: 700;
    border-radius: 8px; padding: 10px 14px; margin-bottom: 8px;
    line-height: 1.4;
  }}
  .footer {{
    background: #f0f4f2; padding: 14px 48px;
    font-size: 11px; color: #8a9399;
    display: flex; justify-content: space-between;
  }}
</style>
</head>
<body>
<div class="page">
  <div class="header">
    <div class="header-left">
      <h1>{doc.get('customer','—')}</h1>
      <div class="use-case">{doc.get('use_case','—')}</div>
      <div class="mdb-badge">MongoDB Proof Point</div>
    </div>
    <div class="header-right">
      {doc['_id']}<br/>
      {datetime.datetime.now().strftime('%Y-%m-%d')}
    </div>
  </div>
  <div class="body">
    <div class="col problem">
      <div class="col-label">Problem</div>
      <h2>The Challenge</h2>
      <p>{prob.get('situation','—')}</p>
      {"<ul>" + neg_outcomes_html + "</ul>" if neg_outcomes_html else ""}
    </div>
    <div class="col solution">
      <div class="col-label">Solution &amp; Why MongoDB</div>
      <h2>What Was Built</h2>
      <p>{sol.get('what_was_built','—')}</p>
      {"<div class='caps'>" + caps_html + "</div>" if caps_html else ""}
      {"<div class='why-mdb'>" + sol.get('why_mongodb','') + "</div>" if sol.get('why_mongodb') else ""}
    </div>
    <div class="col results">
      <div class="col-label">Results</div>
      <h2>Business Outcomes</h2>
      {kpis_html}
      {"<ul style='margin-top:12px'>" + outcomes_html + "</ul>" if outcomes_html else ""}
    </div>
  </div>
  <div class="footer">
    <span>Generated by ACC — Agentic Content Creator</span>
    <span>mongodb.com</span>
  </div>
</div>
</body>
</html>"""

    path = EXPORTS_DIR / f"{_slug(doc)}.html"
    path.write_text(html, encoding="utf-8")
    return path


# ─── PPTX export ────────────────────────────────────────────────────────────

def _export_pptx_file(doc: dict) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prob = doc.get("problem") or {}
    sol  = doc.get("solution") or {}
    res  = doc.get("results") or {}

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slide  = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    def rgb(hex_str):
        h = hex_str.lstrip("#")
        return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    def add_rect(slide, l, t, w, h, fill_hex=None, line_hex=None):
        shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        if fill_hex:
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(fill_hex)
        else:
            shape.fill.background()
        if line_hex:
            shape.line.color.rgb = rgb(line_hex)
            shape.line.width = Emu(12700)
        else:
            shape.line.fill.background()
        return shape

    def add_text(slide, text, l, t, w, h, size, bold=False, color_hex="001E2B",
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
        txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = txBox.text_frame
        tf.word_wrap = wrap
        p  = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = rgb(color_hex)
        return txBox

    # Background
    add_rect(slide, 0, 0, 13.33, 7.5, fill_hex="FFFFFF")

    # Header bar
    add_rect(slide, 0, 0, 13.33, 1.3, fill_hex=MDB_DARK)

    # Header text
    add_text(slide, doc.get("customer", "—"), 0.35, 0.12, 7, 0.55,
             size=26, bold=True, color_hex=MDB_WHITE)
    add_text(slide, doc.get("use_case", "—"), 0.35, 0.68, 7, 0.4,
             size=13, color_hex="9EA8AC")
    add_text(slide, "MongoDB Proof Point", 10.2, 0.2, 2.8, 0.35,
             size=10, bold=True, color_hex=MDB_GREEN, align=PP_ALIGN.RIGHT)
    add_text(slide, doc["_id"] + "  ·  " + datetime.datetime.now().strftime("%Y-%m-%d"),
             10.2, 0.55, 2.8, 0.3, size=9, color_hex="9EA8AC", align=PP_ALIGN.RIGHT)

    # Green accent bar
    add_rect(slide, 0, 1.3, 13.33, 0.04, fill_hex=MDB_GREEN)

    # Three column dividers
    col_w = 4.2
    for i, (label, color_hex, title) in enumerate([
        ("PROBLEM", "D64045", "The Challenge"),
        ("SOLUTION & WHY MONGODB", "0064FF", "What Was Built"),
        ("RESULTS", "00684A", "Business Outcomes"),
    ]):
        x = 0.2 + i * (col_w + 0.17)
        # Column label
        add_text(slide, label, x, 1.48, col_w, 0.25, size=8, bold=True,
                 color_hex=color_hex)
        # Column title
        add_text(slide, title, x, 1.72, col_w, 0.4, size=14, bold=True,
                 color_hex=MDB_DARK)
        # Divider (after first two cols)
        if i < 2:
            add_rect(slide, x + col_w + 0.04, 1.35, 0.02, 5.9, fill_hex="E8EDEB")

    # Problem content
    x0 = 0.2
    prob_txt = (prob.get("situation") or "—")
    add_text(slide, prob_txt, x0, 2.2, col_w, 1.8, size=11, color_hex="3D4F58")
    neg = prob.get("negative_outcomes") or []
    neg_txt = "\n".join(f"• {o}" for o in neg[:4])
    if neg_txt:
        add_text(slide, neg_txt, x0, 4.05, col_w, 2.2, size=10, color_hex="3D4F58")

    # Solution content
    x1 = 0.2 + col_w + 0.17
    sol_txt = sol.get("what_was_built") or "—"
    add_text(slide, sol_txt, x1, 2.2, col_w, 1.5, size=11, color_hex="3D4F58")
    caps = sol.get("mongodb_capabilities") or []
    if caps:
        add_text(slide, " · ".join(caps[:6]), x1, 3.75, col_w, 0.5,
                 size=9, bold=True, color_hex="0064FF")
    why = sol.get("why_mongodb") or ""
    if why:
        add_rect(slide, x1, 4.3, col_w, 0.03, fill_hex=MDB_GREEN)
        add_text(slide, why, x1, 4.38, col_w, 1.8, size=10,
                 italic=True, color_hex="3D4F58")

    # Results content
    x2 = 0.2 + 2 * (col_w + 0.17)
    kpis = (res.get("kpis") or [])[:3]
    y = 2.2
    for kpi in kpis:
        box = add_rect(slide, x2, y, col_w, 0.65, fill_hex=MDB_GREEN)
        add_text(slide, kpi, x2 + 0.12, y + 0.1, col_w - 0.24, 0.5,
                 size=10, bold=True, color_hex=MDB_DARK)
        y += 0.75
    outcomes = res.get("outcomes") or []
    out_txt = "\n".join(f"• {o}" for o in outcomes[:4])
    if out_txt:
        add_text(slide, out_txt, x2, y + 0.1, col_w, 7.2 - y, size=10,
                 color_hex="3D4F58")

    # Footer
    add_rect(slide, 0, 7.22, 13.33, 0.28, fill_hex="F0F4F2")
    add_text(slide, "Generated by ACC — Agentic Content Creator  ·  mongodb.com",
             0.35, 7.24, 12.63, 0.24, size=8, color_hex="8A9399")

    path = EXPORTS_DIR / f"{_slug(doc)}.pptx"
    prs.save(str(path))
    return path


# ─── Tools ──────────────────────────────────────────────────────────────────

@mcp.tool()
def export_html(proof_point_id: str = None) -> str:
    """
    Export the proof point as a clean HTML one-pager (three-column
    Problem / Solution / Results layout, MongoDB-branded). Opens directly
    in a browser. Blocks if required sections are missing.

    Args:
        proof_point_id: Optional. Defaults to the most recently active proof point.
    """
    doc = proof_points.find_one({"_id": proof_point_id}) if proof_point_id else _active()
    if not doc:
        return "❌ No active proof point."

    missing = _missing(doc)
    if missing:
        return (
            f"❌ Export blocked — {doc['_id']} is incomplete.\n"
            f"  Still missing: {'; '.join(missing)}\n"
            f"  Add the missing sections then retry."
        )

    path = _export_html_file(doc)
    proof_points.update_one({"_id": doc["_id"]}, {"$set": {"status": "published"}})

    return (
        f"✅ HTML one-pager exported for **{doc.get('customer')}**.\n"
        f"  File: `{path}`\n"
        f"  Open in browser: file://{path}"
    )


@mcp.tool()
def export_pptx(proof_point_id: str = None) -> str:
    """
    Export the proof point as a single-slide PowerPoint (.pptx) following the
    MongoDB brand style. Three-column layout: Problem / Solution+Why MongoDB /
    Results with KPI tiles. Blocks if required sections are missing.

    Args:
        proof_point_id: Optional. Defaults to the most recently active proof point.
    """
    doc = proof_points.find_one({"_id": proof_point_id}) if proof_point_id else _active()
    if not doc:
        return "❌ No active proof point."

    missing = _missing(doc)
    if missing:
        return (
            f"❌ Export blocked — {doc['_id']} is incomplete.\n"
            f"  Still missing: {'; '.join(missing)}\n"
            f"  Add the missing sections then retry."
        )

    try:
        path = _export_pptx_file(doc)
    except ImportError:
        return "❌ python-pptx not installed. Run: pip install python-pptx"

    proof_points.update_one({"_id": doc["_id"]}, {"$set": {"status": "published"}})

    return (
        f"✅ PowerPoint slide exported for **{doc.get('customer')}**.\n"
        f"  File: `{path}`\n"
        f"  Open with: open '{path}'"
    )


@mcp.tool()
def list_exports() -> str:
    """List all exported files in the exports/ directory."""
    files = sorted(EXPORTS_DIR.glob("PP-*"), key=lambda p: p.stat().st_mtime, reverse=True)
    if not files:
        return "No exports yet."
    lines = [f"**{len(files)} export{'s' if len(files) != 1 else ''}** in `{EXPORTS_DIR}`:\n"]
    for f in files:
        size = f.stat().st_size
        mtime = datetime.datetime.fromtimestamp(f.stat().st_mtime).strftime("%Y-%m-%d %H:%M")
        lines.append(f"- `{f.name}` · {size//1024}KB · {mtime}")
    return "\n".join(lines)


if __name__ == "__main__":
    mcp.run()
